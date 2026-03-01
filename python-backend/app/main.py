from fastapi import FastAPI, HTTPException, Query
from pydantic import BaseModel
from pathlib import Path
from dotenv import load_dotenv
from typing import Optional, List
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Boolean, Float, Text, ForeignKey, inspect, func, event, text
from sqlalchemy.orm import declarative_base, Session, relationship
import os
import datetime
import uuid
import hashlib
import mimetypes

load_dotenv()

app = FastAPI(title="DocRack AI Engine")

# ---------------------------------------------------------------------------
# In-memory engine registry — set by POST /init-db, cleared on app restart.
# All DB-dependent routes must check this before operating.
# ---------------------------------------------------------------------------
active_engine = None

# Accepted theme values — validated before any DB write.
_VALID_THEMES = {"light", "dark"}

# Allowed file extensions for registration.
_ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".csv", ".png", ".jpg", ".jpeg"}

# Max files per registration request.
_MAX_FILES_PER_REQUEST = 50

# Max characters stored in extracted_text column.
_MAX_EXTRACTED_TEXT_LENGTH = 5000

# ---------------------------------------------------------------------------
# SQLAlchemy ORM — schema defined and owned exclusively by this Python backend.
# No other layer may create, alter, or drop these tables.
# ---------------------------------------------------------------------------
Base = declarative_base()


class UserMeta(Base):
    __tablename__ = "user_meta"

    id = Column(Integer, primary_key=True, autoincrement=True)
    mongo_user_id = Column(String, nullable=False, unique=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)


class Settings(Base):
    __tablename__ = "settings"

    key = Column(String, primary_key=True)
    value = Column(String, nullable=True)


# ---------------------------------------------------------------------------
# Smart DataRoom models
# ---------------------------------------------------------------------------

def _generate_uuid():
    return str(uuid.uuid4())


class DataRoom(Base):
    __tablename__ = "datarooms"

    id = Column(String, primary_key=True, default=_generate_uuid)
    name = Column(String, nullable=False)
    description = Column(Text, nullable=True)
    created_by_ai = Column(Boolean, default=False)
    status = Column(String, default="active")  # active | archived
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)

    folders = relationship("Folder", back_populates="dataroom", cascade="all, delete-orphan")
    files = relationship("File", back_populates="dataroom", cascade="all, delete-orphan")


class Folder(Base):
    __tablename__ = "folders"

    id = Column(String, primary_key=True, default=_generate_uuid)
    dataroom_id = Column(String, ForeignKey("datarooms.id", ondelete="CASCADE"), nullable=False)
    name = Column(String, nullable=False)
    context = Column(Text, nullable=False)
    parent_id = Column(String, ForeignKey("folders.id"), nullable=True)
    display_order = Column(Integer, default=0)
    created_by_ai = Column(Boolean, default=False)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)

    dataroom = relationship("DataRoom", back_populates="folders")
    children = relationship("Folder", backref="parent", remote_side=[id])
    files = relationship("File", back_populates="folder")
    classifications = relationship("Classification", back_populates="folder", cascade="all, delete-orphan")


class File(Base):
    __tablename__ = "files"

    id = Column(String, primary_key=True, default=_generate_uuid)
    dataroom_id = Column(String, ForeignKey("datarooms.id", ondelete="CASCADE"), nullable=False)
    folder_id = Column(String, ForeignKey("folders.id", ondelete="SET NULL"), nullable=True)
    original_name = Column(String, nullable=False)
    original_path = Column(Text, nullable=False)
    file_extension = Column(String, nullable=False)
    mime_type = Column(String, nullable=True)
    size_bytes = Column(Integer, nullable=True)
    checksum = Column(String, nullable=True)
    extracted_text = Column(Text, nullable=True)
    status = Column(String, default="registered")  # registered | processing | classified | error
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)

    dataroom = relationship("DataRoom", back_populates="files")
    folder = relationship("Folder", back_populates="files")
    classifications = relationship("Classification", back_populates="file", cascade="all, delete-orphan")


class Classification(Base):
    __tablename__ = "classifications"

    id = Column(String, primary_key=True, default=_generate_uuid)
    file_id = Column(String, ForeignKey("files.id", ondelete="CASCADE"), nullable=False)
    folder_id = Column(String, ForeignKey("folders.id", ondelete="CASCADE"), nullable=False)
    confidence = Column(Float, nullable=True)
    reasoning = Column(Text, nullable=True)
    classified_at = Column(DateTime, default=datetime.datetime.utcnow)

    file = relationship("File", back_populates="classifications")
    folder = relationship("Folder", back_populates="classifications")


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _require_db():
    """
    Returns active_engine or raises 503 if /init-db has not been called yet.
    Call this at the start of every DB-dependent route.
    """
    if active_engine is None:
        raise HTTPException(
            status_code=503,
            detail="Database not initialised. Call POST /init-db first.",
        )
    return active_engine


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def _compute_checksum(file_path: str) -> str:
    """Compute SHA-256 checksum of a file."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def _extract_text_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz
    text_parts = []
    doc = fitz.open(file_path)
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n".join(text_parts)


def _extract_text_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs)


def _extract_text_xlsx(file_path: str) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        text_parts.append(f"[Sheet: {sheet_name}]")
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(cell) for cell in row if cell is not None]
            if cell_values:
                text_parts.append(" | ".join(cell_values))
    wb.close()
    return "\n".join(text_parts)


def _extract_text_pptx(file_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
    return "\n".join(text_parts)


def _extract_text_plain(file_path: str) -> str:
    """Extract text from TXT/CSV with encoding fallback."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read()


def _extract_text(file_path: str, extension: str, filename: str) -> str:
    """
    Extract text from a file based on its extension.
    Returns the extracted text or raises an exception on failure.
    """
    ext = extension.lower()

    if ext == ".pdf":
        return _extract_text_pdf(file_path)
    elif ext == ".docx":
        return _extract_text_docx(file_path)
    elif ext == ".xlsx":
        return _extract_text_xlsx(file_path)
    elif ext == ".pptx":
        return _extract_text_pptx(file_path)
    elif ext in (".txt", ".csv"):
        return _extract_text_plain(file_path)
    elif ext in (".png", ".jpg", ".jpeg"):
        return f"[Image: {filename}]"
    else:
        return ""


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.get("/health")
def health():
    result = {"status": "ok"}
    if active_engine is not None:
        inspector = inspect(active_engine)
        result["tables"] = inspector.get_table_names()
    return result


# ---- Database initialisation -----------------------------------------------

class InitDbRequest(BaseModel):
    database_path: str
    mongo_user_id: str


@app.post("/init-db")
def init_db(request: InitDbRequest):
    global active_engine

    db_path = request.database_path

    # Security: only absolute paths are accepted
    if not os.path.isabs(db_path):
        raise HTTPException(
            status_code=400,
            detail="database_path must be an absolute path",
        )

    # Security: reject any path containing '..' components
    path_obj = Path(db_path)
    if ".." in path_obj.parts:
        raise HTTPException(
            status_code=400,
            detail="Path traversal is not allowed in database_path",
        )

    # Resolve to normalise redundant separators without changing the target
    try:
        resolved = str(path_obj.resolve())
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid database_path")

    # Create parent directory if it does not exist
    parent_dir = os.path.dirname(resolved)
    try:
        os.makedirs(parent_dir, exist_ok=True)
    except OSError as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not create database directory: {exc}",
        )

    # Create engine and initialise schema — Python owns all DDL
    engine = create_engine(f"sqlite:///{resolved}", echo=False)

    # SQLite requires an explicit pragma to enforce foreign key constraints.
    # Without this, ON DELETE CASCADE / SET NULL are silently ignored.
    @event.listens_for(engine, "connect")
    def _set_sqlite_pragma(dbapi_conn, _connection_record):
        cursor = dbapi_conn.cursor()
        cursor.execute("PRAGMA foreign_keys = ON")
        cursor.close()

    # Schema migration: if the files table exists with stale columns
    # (stored_name, local_path from old schema), drop files + classifications
    # so create_all rebuilds them with the current ORM definitions.
    with engine.connect() as conn:
        result = conn.execute(text(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='files'"
        ))
        if result.fetchone():
            columns = [
                row[1] for row in conn.execute(text("PRAGMA table_info(files)"))
            ]
            if "original_path" not in columns:
                conn.execute(text("DROP TABLE IF EXISTS classifications"))
                conn.execute(text("DROP TABLE IF EXISTS files"))
                conn.commit()

    Base.metadata.create_all(engine)

    # Upsert user_meta row — idempotent, one row per mongo_user_id
    with Session(engine) as session:
        existing = (
            session.query(UserMeta)
            .filter_by(mongo_user_id=request.mongo_user_id)
            .first()
        )
        if not existing:
            session.add(UserMeta(mongo_user_id=request.mongo_user_id))
            session.commit()

    # Register engine only after schema + seed succeed
    active_engine = engine

    return {
        "status": "success",
        "message": "Database initialized",
        "path": resolved,
    }


# ---- Theme settings ---------------------------------------------------------

@app.get("/settings/theme")
def get_theme():
    """
    Returns the stored theme for the active user.
    Defaults to "light" if no theme has been persisted yet.
    Requires /init-db to have been called first.
    """
    engine = _require_db()
    with Session(engine) as session:
        row = session.query(Settings).filter_by(key="theme").first()
    return {"theme": row.value if row else "light"}


class ThemeRequest(BaseModel):
    theme: str


@app.post("/settings/theme")
def set_theme(request: ThemeRequest):
    """
    Persists the theme to the settings table.
    Only "light" and "dark" are accepted — all other values are rejected.
    Upserts the existing row so there is never more than one theme entry.
    Requires /init-db to have been called first.
    """
    if request.theme not in _VALID_THEMES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid theme value '{request.theme}'. Allowed: light, dark.",
        )

    engine = _require_db()
    with Session(engine) as session:
        row = session.query(Settings).filter_by(key="theme").first()
        if row:
            row.value = request.theme
        else:
            session.add(Settings(key="theme", value=request.theme))
        session.commit()

    return {"status": "success", "theme": request.theme}


# ---- DataRoom & Folder CRUD ------------------------------------------------

# -- Pydantic request models --

class CreateDataRoomRequest(BaseModel):
    name: str
    description: Optional[str] = None


class UpdateDataRoomRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None


class CreateFolderRequest(BaseModel):
    name: str
    context: str
    parent_id: Optional[str] = None


class UpdateFolderRequest(BaseModel):
    name: Optional[str] = None
    context: Optional[str] = None
    parent_id: Optional[str] = "__unset__"  # distinguish between "not provided" and "set to null"


# -- Pydantic request models for Files --

class RegisterFilesRequest(BaseModel):
    dataroom_id: str
    file_paths: List[str]


class RelocateFileRequest(BaseModel):
    new_path: str


class MoveToFolderRequest(BaseModel):
    folder_id: Optional[str] = None


class RenameFileRequest(BaseModel):
    new_name: str


# -- Serialisation helpers --

def _dt(val):
    """Convert a datetime to ISO-8601 string, or None."""
    return val.isoformat() if val else None


def _dataroom_dict(dr):
    return {
        "id": dr.id,
        "name": dr.name,
        "description": dr.description,
        "created_by_ai": dr.created_by_ai,
        "status": dr.status,
        "created_at": _dt(dr.created_at),
        "updated_at": _dt(dr.updated_at),
    }


def _folder_dict(f):
    return {
        "id": f.id,
        "dataroom_id": f.dataroom_id,
        "name": f.name,
        "context": f.context,
        "parent_id": f.parent_id,
        "display_order": f.display_order,
        "created_by_ai": f.created_by_ai,
        "created_at": _dt(f.created_at),
        "updated_at": _dt(f.updated_at),
    }


def _file_dict(f):
    return {
        "id": f.id,
        "dataroom_id": f.dataroom_id,
        "folder_id": f.folder_id,
        "original_name": f.original_name,
        "original_path": f.original_path,
        "file_extension": f.file_extension,
        "mime_type": f.mime_type,
        "size_bytes": f.size_bytes,
        "checksum": f.checksum,
        "status": f.status,
        "created_at": _dt(f.created_at),
        "updated_at": _dt(f.updated_at),
    }


# -- DataRoom endpoints --

@app.post("/datarooms", status_code=201)
def create_dataroom(request: CreateDataRoomRequest):
    if not request.name or not request.name.strip():
        raise HTTPException(status_code=400, detail="DataRoom name is required.")

    engine = _require_db()
    dr = DataRoom(name=request.name.strip(), description=request.description)

    with Session(engine) as session:
        session.add(dr)
        session.commit()
        session.refresh(dr)
        return _dataroom_dict(dr)


@app.get("/datarooms")
def list_datarooms():
    engine = _require_db()
    with Session(engine) as session:
        rows = session.query(DataRoom).order_by(DataRoom.created_at.desc()).all()

        results = []
        for dr in rows:
            folder_count = session.query(func.count(Folder.id)).filter(Folder.dataroom_id == dr.id).scalar()
            file_count = session.query(func.count(File.id)).filter(File.dataroom_id == dr.id).scalar()
            d = _dataroom_dict(dr)
            d["folder_count"] = folder_count
            d["file_count"] = file_count
            results.append(d)

        return results


@app.get("/datarooms/{dataroom_id}")
def get_dataroom(dataroom_id: str):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        folders_raw = session.query(Folder).filter_by(dataroom_id=dr.id).order_by(Folder.display_order).all()
        folders = []
        for f in folders_raw:
            fd = _folder_dict(f)
            fd["file_count"] = session.query(func.count(File.id)).filter(File.folder_id == f.id).scalar()
            folders.append(fd)

        files = [
            _file_dict(f) for f in
            session.query(File).filter_by(dataroom_id=dr.id).order_by(File.created_at.desc()).all()
        ]

        result = _dataroom_dict(dr)
        result["folders"] = folders
        result["files"] = files
        return result


@app.put("/datarooms/{dataroom_id}")
def update_dataroom(dataroom_id: str, request: UpdateDataRoomRequest):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        if request.name is not None:
            if not request.name.strip():
                raise HTTPException(status_code=400, detail="DataRoom name cannot be empty.")
            dr.name = request.name.strip()

        if request.description is not None:
            dr.description = request.description

        dr.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(dr)
        return _dataroom_dict(dr)


@app.delete("/datarooms/{dataroom_id}")
def delete_dataroom(dataroom_id: str):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        session.delete(dr)
        session.commit()
        return {"success": True, "deleted_id": dataroom_id}


# -- Folder endpoints --

@app.post("/datarooms/{dataroom_id}/folders", status_code=201)
def create_folder(dataroom_id: str, request: CreateFolderRequest):
    if not request.name or not request.name.strip():
        raise HTTPException(status_code=400, detail="Folder name is required.")
    if not request.context or not request.context.strip():
        raise HTTPException(status_code=400, detail="Folder context is required.")

    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        # Validate parent folder if provided
        if request.parent_id is not None:
            parent = session.query(Folder).filter_by(id=request.parent_id).first()
            if not parent:
                raise HTTPException(status_code=404, detail="Parent folder not found.")
            if parent.dataroom_id != dataroom_id:
                raise HTTPException(status_code=400, detail="Parent folder does not belong to this DataRoom.")

        folder = Folder(
            dataroom_id=dataroom_id,
            name=request.name.strip(),
            context=request.context.strip(),
            parent_id=request.parent_id,
        )
        session.add(folder)
        session.commit()
        session.refresh(folder)
        return _folder_dict(folder)


@app.get("/datarooms/{dataroom_id}/folders")
def list_folders(dataroom_id: str):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        folders_raw = session.query(Folder).filter_by(dataroom_id=dataroom_id).order_by(Folder.display_order).all()
        results = []
        for f in folders_raw:
            fd = _folder_dict(f)
            fd["file_count"] = session.query(func.count(File.id)).filter(File.folder_id == f.id).scalar()
            results.append(fd)

        return results


@app.put("/folders/{folder_id}")
def update_folder(folder_id: str, request: UpdateFolderRequest):
    engine = _require_db()
    with Session(engine) as session:
        folder = session.query(Folder).filter_by(id=folder_id).first()
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found.")

        if request.name is not None:
            if not request.name.strip():
                raise HTTPException(status_code=400, detail="Folder name cannot be empty.")
            folder.name = request.name.strip()

        if request.context is not None:
            if not request.context.strip():
                raise HTTPException(status_code=400, detail="Folder context cannot be empty.")
            folder.context = request.context.strip()

        # Handle parent_id change (folder move)
        if request.parent_id != "__unset__":
            if request.parent_id is not None:
                # Validate parent exists and belongs to same DataRoom
                parent = session.query(Folder).filter_by(id=request.parent_id).first()
                if not parent:
                    raise HTTPException(status_code=404, detail="Parent folder not found.")
                if parent.dataroom_id != folder.dataroom_id:
                    raise HTTPException(status_code=400, detail="Parent folder does not belong to the same DataRoom.")
                # Prevent circular reference — folder cannot be its own ancestor
                if request.parent_id == folder_id:
                    raise HTTPException(status_code=400, detail="A folder cannot be its own parent.")
            folder.parent_id = request.parent_id

        folder.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(folder)
        return _folder_dict(folder)


@app.delete("/folders/{folder_id}")
def delete_folder(folder_id: str):
    engine = _require_db()
    with Session(engine) as session:
        folder = session.query(Folder).filter_by(id=folder_id).first()
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found.")

        # Files in this folder get folder_id = NULL (handled by ON DELETE SET NULL),
        # but SQLAlchemy relationship cache may be stale — flush explicitly.
        session.query(File).filter_by(folder_id=folder_id).update({"folder_id": None})
        session.delete(folder)
        session.commit()
        return {"success": True}


# ---- File endpoints ---------------------------------------------------------

@app.post("/files/register", status_code=201)
def register_files(request: RegisterFilesRequest):
    """
    Register file paths and extract text content.
    Files stay at their original location — only the path is stored in SQLite.
    """
    engine = _require_db()

    # Validate file count limit
    if len(request.file_paths) > _MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Maximum {_MAX_FILES_PER_REQUEST} files per request. Received {len(request.file_paths)}.",
        )

    if not request.file_paths:
        raise HTTPException(status_code=400, detail="file_paths list cannot be empty.")

    with Session(engine) as session:
        # Verify DataRoom exists
        dr = session.query(DataRoom).filter_by(id=request.dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        # Get all existing paths in this DataRoom for duplicate detection
        existing_paths = set(
            row[0] for row in
            session.query(File.original_path).filter_by(dataroom_id=request.dataroom_id).all()
        )

        registered = []
        rejected = []

        for file_path in request.file_paths:
            # Normalise path
            normalised = os.path.normpath(file_path)

            # Check: must be absolute
            if not os.path.isabs(normalised):
                rejected.append({"path": file_path, "reason": "Path must be absolute."})
                continue

            # Check: file extension
            _, ext = os.path.splitext(normalised)
            ext_lower = ext.lower()
            if ext_lower not in _ALLOWED_EXTENSIONS:
                rejected.append({"path": file_path, "reason": f"Extension '{ext_lower}' is not allowed. Allowed: {', '.join(sorted(_ALLOWED_EXTENSIONS))}"})
                continue

            # Check: file exists on disk
            if not os.path.exists(normalised):
                rejected.append({"path": file_path, "reason": "File does not exist at this path."})
                continue

            if not os.path.isfile(normalised):
                rejected.append({"path": file_path, "reason": "Path is not a file."})
                continue

            # Check: duplicate within this DataRoom
            if normalised in existing_paths:
                rejected.append({"path": file_path, "reason": "File already registered in this DataRoom."})
                continue

            # Read metadata
            original_name = os.path.basename(normalised)
            try:
                size_bytes = os.path.getsize(normalised)
            except OSError as exc:
                rejected.append({"path": file_path, "reason": f"Cannot read file size: {exc}"})
                continue

            mime_type, _ = mimetypes.guess_type(normalised)

            # Compute checksum
            try:
                checksum = _compute_checksum(normalised)
            except (OSError, PermissionError) as exc:
                rejected.append({"path": file_path, "reason": f"Cannot read file for checksum: {exc}"})
                continue

            # Extract text
            try:
                extracted = _extract_text(normalised, ext_lower, original_name)
            except Exception as exc:
                # File is corrupted or locked — register with error status
                file_record = File(
                    dataroom_id=request.dataroom_id,
                    original_name=original_name,
                    original_path=normalised,
                    file_extension=ext_lower,
                    mime_type=mime_type,
                    size_bytes=size_bytes,
                    checksum=checksum,
                    extracted_text=f"[Extraction error: {exc}]",
                    status="error",
                )
                session.add(file_record)
                session.flush()
                existing_paths.add(normalised)
                rejected.append({"path": file_path, "reason": f"Text extraction failed: {exc}"})
                continue

            # Truncate extracted text to max length
            if extracted and len(extracted) > _MAX_EXTRACTED_TEXT_LENGTH:
                extracted = extracted[:_MAX_EXTRACTED_TEXT_LENGTH]

            # Create file record
            file_record = File(
                dataroom_id=request.dataroom_id,
                original_name=original_name,
                original_path=normalised,
                file_extension=ext_lower,
                mime_type=mime_type,
                size_bytes=size_bytes,
                checksum=checksum,
                extracted_text=extracted,
                status="registered",
            )
            session.add(file_record)
            session.flush()
            session.refresh(file_record)
            existing_paths.add(normalised)

            registered.append(_file_dict(file_record))

        session.commit()

    return {
        "registered": registered,
        "rejected": rejected,
        "total_registered": len(registered),
        "total_rejected": len(rejected),
    }


@app.get("/files/{file_id}")
def get_file(file_id: str):
    """Returns full file metadata including original_path."""
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        result = _file_dict(file_record)
        result["extracted_text"] = file_record.extracted_text
        return result


@app.post("/files/{file_id}/check-exists")
def check_file_exists(file_id: str):
    """Checks if the file still exists at its original_path on disk."""
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        exists = os.path.exists(file_record.original_path) and os.path.isfile(file_record.original_path)
        return {"exists": exists, "path": file_record.original_path}


@app.put("/files/{file_id}/relocate")
def relocate_file(file_id: str, request: RelocateFileRequest):
    """
    Update the stored path after a user has moved a file on disk.
    Validates the new path exists and checksum matches.
    """
    engine = _require_db()

    new_path = os.path.normpath(request.new_path)

    if not os.path.isabs(new_path):
        raise HTTPException(status_code=400, detail="new_path must be an absolute path.")

    if not os.path.exists(new_path) or not os.path.isfile(new_path):
        raise HTTPException(status_code=400, detail="File does not exist at the new path.")

    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        # Verify checksum matches
        if file_record.checksum:
            try:
                new_checksum = _compute_checksum(new_path)
            except (OSError, PermissionError) as exc:
                raise HTTPException(status_code=400, detail=f"Cannot read file at new path: {exc}")

            if new_checksum != file_record.checksum:
                raise HTTPException(
                    status_code=400,
                    detail="Checksum mismatch. The file at the new path does not match the original file.",
                )

        file_record.original_path = new_path
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(file_record)
        return _file_dict(file_record)


@app.put("/files/{file_id}/move-to-folder")
def move_to_folder(file_id: str, request: MoveToFolderRequest):
    """Move file to a different virtual folder (null = unclassified)."""
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        if request.folder_id is not None:
            folder = session.query(Folder).filter_by(id=request.folder_id).first()
            if not folder:
                raise HTTPException(status_code=404, detail="Folder not found.")
            if folder.dataroom_id != file_record.dataroom_id:
                raise HTTPException(status_code=400, detail="Folder does not belong to the same DataRoom.")

        file_record.folder_id = request.folder_id
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(file_record)
        return _file_dict(file_record)


@app.delete("/files/{file_id}")
def delete_file(file_id: str, delete_from_system: bool = Query(default=False)):
    """
    Remove file record from SQLite.
    If delete_from_system=true, also deletes the actual file from disk.
    """
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        deleted_from_system = False
        if delete_from_system:
            try:
                if os.path.exists(file_record.original_path):
                    os.remove(file_record.original_path)
                    deleted_from_system = True
            except OSError as exc:
                raise HTTPException(
                    status_code=500,
                    detail=f"Failed to delete file from disk: {exc}",
                )

        session.delete(file_record)
        session.commit()
        return {"success": True, "deleted_from_system": deleted_from_system}


@app.get("/datarooms/{dataroom_id}/files")
def list_dataroom_files(
    dataroom_id: str,
    folder_id: Optional[str] = Query(default=None),
    include_subfolders: bool = Query(default=False),
    status: Optional[str] = Query(default=None),
):
    """
    List files in a DataRoom with optional filters.
    If folder_id is provided, filters to that folder.
    If include_subfolders is true, also includes files from nested child folders.
    """
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        query = session.query(File).filter(File.dataroom_id == dataroom_id)

        if folder_id is not None:
            if include_subfolders:
                # Collect folder_id and all descendant folder IDs
                target_folder_ids = _collect_subfolder_ids(session, folder_id)
                query = query.filter(File.folder_id.in_(target_folder_ids))
            else:
                query = query.filter(File.folder_id == folder_id)

        if status is not None:
            query = query.filter(File.status == status)

        files = query.order_by(File.created_at.desc()).all()

        results = []
        for f in files:
            fd = _file_dict(f)
            if f.folder:
                fd["folder_name"] = f.folder.name
            else:
                fd["folder_name"] = None
            results.append(fd)

        return results


@app.put("/files/{file_id}/rename")
def rename_file(file_id: str, request: RenameFileRequest):
    """
    Update display name in SQLite only — does NOT rename the actual file on disk.
    """
    engine = _require_db()

    if not request.new_name or not request.new_name.strip():
        raise HTTPException(status_code=400, detail="new_name cannot be empty.")

    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        file_record.original_name = request.new_name.strip()
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(file_record)
        return _file_dict(file_record)


# ---- AI Classification endpoints --------------------------------------------

class ClassifyRequest(BaseModel):
    dataroom_id: str
    file_ids: List[str]


class GenerateDataRoomRequest(BaseModel):
    dataroom_name: str
    dataroom_description: Optional[str] = None
    file_ids: List[str]


@app.post("/ai/classify")
async def ai_classify(request: ClassifyRequest):
    """
    Classify files into existing DataRoom folders using Gemini AI.
    Files are matched to folders based on content and folder context.
    """
    from app.services.classification_service import classify_files

    engine = _require_db()

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    if len(request.file_ids) > _MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Maximum {_MAX_FILES_PER_REQUEST} files per request. Received {len(request.file_ids)}.",
        )

    # Verify DataRoom exists
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=request.dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

    try:
        result = await classify_files(engine, request.dataroom_id, request.file_ids)
        return result
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except RuntimeError as e:
        raise HTTPException(status_code=502, detail=str(e))


@app.post("/ai/generate-dataroom")
async def ai_generate_dataroom(request: GenerateDataRoomRequest):
    """
    Generate a new DataRoom with AI-created folder structure and file assignments.
    Uses Gemini to analyze files and create an organized structure.
    """
    from app.services.classification_service import generate_dataroom

    engine = _require_db()

    if not request.dataroom_name or not request.dataroom_name.strip():
        raise HTTPException(status_code=400, detail="dataroom_name is required.")

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    if len(request.file_ids) > _MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Maximum {_MAX_FILES_PER_REQUEST} files per request. Received {len(request.file_ids)}.",
        )

    try:
        result = await generate_dataroom(
            engine,
            request.dataroom_name.strip(),
            request.dataroom_description,
            request.file_ids,
        )
        return result
    except RuntimeError as e:
        raise HTTPException(status_code=502, detail=str(e))


# ---------------------------------------------------------------------------
# Subfolder traversal helper
# ---------------------------------------------------------------------------

def _collect_subfolder_ids(session, folder_id: str) -> List[str]:
    """
    BFS traversal to collect a folder and all its descendant folder IDs.
    """
    result = [folder_id]
    queue = [folder_id]
    while queue:
        current_id = queue.pop(0)
        children = session.query(Folder.id).filter_by(parent_id=current_id).all()
        for (child_id,) in children:
            result.append(child_id)
            queue.append(child_id)
    return result
